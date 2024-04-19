from pptx import Presentation
from pptx.slide import Slide
from argparse import ArgumentParser
from pkg_resources import get_distribution
from cast_arg.config import Config

from pandas import DataFrame

from copy import deepcopy

__author__ = "Nevin Kaplan"
__email__ = "n.kaplan@castsoftware.com"
__copyright__ = "Copyright 2023, CAST Software"

def duplicate_slide(ppt: Presentation, source: Slide):
    """
    Duplicate the slide with the given number in presentation.
    Adds the new slide by default at the end of the presentation.

    :param ppt:
    :param slide_index: Slide number
    :return:
    """
    #source = ppt2.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy existing references of known type
    # e.g. hyperlinks
    known_refs = [
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    ]
    for rel in _object_rels(source.part):
        if rel.reltype in known_refs:
            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    return dest

def delete_slide(prs,idx=-1):
    if idx < 0:
        rng = range(len(prs.slides)-1, -1, -1)
    else:
        rng = range(idx,idx+1)

    for i in rng: 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    layout_0 = pres.slide_layouts[blank_layout_id]
    for shape in layout_0.shapes:
        sp = shape.element
        sp.getparent().remove(sp)
    return layout_0

def _object_rels(obj):
    try:
        rels = obj.rels

        # Change required for python-pptx 0.6.22
        check_rels_content = [k for k in rels]
        if isinstance(check_rels_content.pop(), str):
            return [v for k, v in rels.items()]
        else:
            return [k for k in rels]
    except:
        return []


#def duplicate_slide(ppt: Presentation,ppt2: Presentation, slide_index: int):

def _exp_add_slide(ppt, slide_layout):
    """
    Function to handle slide creation in the Presentation, to avoid issues caused by default implementation.

    :param slide_layout:
    :return:
    """

    def generate_slide_partname(self):
        """Return |PackURI| instance containing next available slide partname."""
        from pptx.opc.packuri import PackURI

        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [k.target_partname for k in _object_rels(self)]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            import random
            import string

            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (
                random_part,
                len(sldIdLst) + 1,
            )

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """
        Return an (rId, slide) pair of a newly created blank slide that
        inherits appearance from *slide_layout*.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.slide import SlidePart

        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)

def remove_shape(shape):
    """
    Helper to remove a specific shape.

    :source: https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

    :param shape:
    :return:
    """
    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree

def copy_shapes(source, dest):
    """
    Helper to copy shapes handling edge cases.

    :param source:
    :param dest:
    :return:
    """
    from pptx.shapes.group import GroupShape
    import copy

    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            import io

            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            result = clone_chart(shape, dest)
        else:
            import copy

            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]

def clone_chart(graphical_frame, dest):
    """
    Helper to clone a chart with related styling.

    :param graphical_frame: General shape containing the .chart property
    :param dest: Shapes object on which to add the new chart
    :return:
    """
    chart = graphical_frame.chart

    df = chart_to_dataframe(graphical_frame)
    chart_data = dataframe_to_chart_data(df)

    new_chart = dest.shapes.add_chart(
        chart.chart_type,
        graphical_frame.left,
        graphical_frame.top,
        graphical_frame.width,
        graphical_frame.height,
        chart_data,
    )

def dataframe_to_chart_data(df):
    """
    Transforms a DataFrame to a CategoryChartData for PPT compilation.

    The indexes of the DataFrame are the categories, with each column becoming a series.

    :param df:
    :return:
    """
    from pptx.chart.data import CategoryChartData
    import numpy as np

    copy_data = CategoryChartData()
    copy_data.categories = df.index.astype(str).to_list()

    edge_cases = 0
    for c in df.columns:
        series_data = df[c].copy()
        fixed_series_data = series_data.replace([np.inf, -np.inf, np.nan], None)

        edge_cases = edge_cases + np.count_nonzero(fixed_series_data != series_data)

        copy_data.add_series(str(c), fixed_series_data.to_list())

    # Warning over data filled for compatibility
    if edge_cases > 0:
        import warnings

        warnings.warn("Series data containing NaN/INF values: filled to empty")

    return copy_data

def chart_to_dataframe(graphical_frame) -> DataFrame:
    """
    Helper to parse chart data to a DataFrame.

    :source: https://openpyxl.readthedocs.io/en/stable/pandas.html

    :param graphical_frame:
    :return:
    """
    from openpyxl import load_workbook

    from io import BytesIO

    wb = load_workbook(
        BytesIO(graphical_frame.chart.part.chart_workbook.xlsx_part.blob),
        read_only=True,
    )

    ws = wb.active

    from itertools import islice
    import pandas as pd

    data = ws.values
    cols = next(data)[1:]
    data = list(data)
    idx = [r[0] for r in data]
    data = (islice(r, 1, None) for r in data)
    df = pd.DataFrame(data, index=idx, columns=cols)

    return df


def clone_deck(prs1, prs2):
    """Duplicate each slide in prs2 and "moves" it into prs1.
    Adds slides to the end of the presentation"""

    for slide in prs2.slides:
        sl = prs1.slides.add_slide(_get_blank_slide_layout(prs1))
        for shape in slide.shapes:
            newel = deepcopy(shape.element)
            sl.shapes._spTree.insert_element_before(newel, 'p:extLst')
        try:
            sl.shapes.title.text = slide.shapes.title.text
            sl.placeholders[0].text = slide.placeholders[0].text
        except Exception as e:
            print(f"Error \"{e}\", suppressing it...")
    return prs1

from pptx.enum.shapes import MSO_SHAPE_TYPE

def _get_shape_by_name(prs, name, use_slide=None, all=False):
    rslt = []

    slides = prs.slides
    if use_slide != None:
        slides = [use_slide] 

    for slide in slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grp_shape in shape.shapes:
                    parts = grp_shape.name.split(':')
                    shape_name = parts[0]
                    if shape_name == name:
                        if not all:
                            return grp_shape
                        else:
                            rslt.append(grp_shape)
            else:
                parts = shape.name.split(':')
                shape_name = parts[0]
                if shape_name == name:
                    if not all:
                        return shape
                    else:
                        rslt.append(shape)
    if len(rslt) == 0:
        return None
    else: 
        return rslt

def move_slide(ppt: Presentation,old_index, new_index):
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

if __name__ == '__main__':
    version = get_distribution('com.castsoftware.uc.arg').version
    print(f'\nCAST Assessment Deck Generation Tool (ARG), v{version}')
    print(f'com.castsoftware.uc.python.common v{get_distribution("com.castsoftware.uc.python.common").version}')
    print('Copyright (c) 2024 CAST Software Inc.')
    print('If you need assistance, please contact oneclick@castsoftware.com')

    parser = ArgumentParser(description='Assessment Report Generation Tool')
    parser.add_argument('-c','--config', required=True, help='Configuration properties file')

    args = parser.parse_args()
    config=Config(args.config)

    template = Presentation(config.template)
#    output = Presentation()

    duplicate_slide(template,template.slides[1])
    move_slide(template,len(template.slides)-1,1)
#    output.save('test.pptx')
    template.save('test.pptx')

    pass


