from __future__ import print_function

import collections
import collections.abc
import pptx
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Length

import numpy as np

from PIL import Image
import os
from rapidfuzz import process as fuze_process
from operator import attrgetter

from tqdm import tqdm
from pptx2md.global_var import g
from pptx2md import global_var

import pptx2md.outputter as outputter

from pptx2md.columns import is_two_column_text, assign_shapes

from pptx2md.utils_optim import normal_pdf, fit_column_model


picture_count = 0

global out


# pptx type defination rules
def is_title(shape):
  if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
                               or shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
                               or shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE
                               or shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
    return True
  return False


def is_text_block(shape):
  if shape.has_text_frame:
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
      return True
    if len(shape.text) > g.text_block_threshold:
      return True
  return False


def is_list_block(shape):
  levels = []
  for para in shape.text_frame.paragraphs:
    if para.level not in levels:
      levels.append(para.level)
    if para.level != 0 or len(levels) > 1:
      return True
  return False


def is_accent(font):
  if font.underline or font.italic or (
      font.color.type == MSO_COLOR_TYPE.SCHEME and
      (font.color.theme_color == MSO_THEME_COLOR.ACCENT_1 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_2
       or font.color.theme_color == MSO_THEME_COLOR.ACCENT_3 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_4
       or font.color.theme_color == MSO_THEME_COLOR.ACCENT_5 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_6)):
    return True
  return False


def is_strong(font):
  if font.bold or (font.color.type == MSO_COLOR_TYPE.SCHEME and (font.color.theme_color == MSO_THEME_COLOR.DARK_1
                                                                 or font.color.theme_color == MSO_THEME_COLOR.DARK_2)):
    return True
  return False


def get_formatted_text(para):
  res = ''
  for run in para.runs:
    text = run.text
    if text == '':
      continue
    if not g.disable_escaping:
      text = out.get_escaped(text)
    try:
      if run.hyperlink.address:
        text = out.get_hyperlink(text, run.hyperlink.address)
    except:
      text = out.get_hyperlink(text, 'error:ppt-link-parsing-issue')
    if is_accent(run.font):
      text = out.get_accent(text)
    elif is_strong(run.font):
      text = out.get_strong(text)
    if not g.disable_color:
      if run.font.color.type == MSO_COLOR_TYPE.RGB:
        text = out.get_colored(text, run.font.color.rgb)
    res += text
  return res.strip()


def process_title(shape, slide_idx):
  global out
  notes = []
  text = shape.text_frame.text.strip()
  if g.use_custom_title:
    res = fuze_process.extractOne(text, g.titles.keys(), score_cutoff=92)
    if not res:
      g.max_custom_title
      out.put_title(text, g.max_custom_title + 1)
    else:
      notes.append(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
      out.put_title(res[0], g.titles[res[0]])
  else:
    out.put_title(text, 1)

  return notes


def process_text_block(shape, _):
  global out
  if is_list_block(shape):
    # generate list block
    for para in shape.text_frame.paragraphs:
      if para.text.strip() == '':
        continue
      text = get_formatted_text(para)
      out.put_list(text, para.level)
    out.write('\n')
  else:
    # generate paragraph block
    for para in shape.text_frame.paragraphs:
      if para.text.strip() == '':
        continue
      text = get_formatted_text(para)
      out.put_para(text)
  return []


def process_notes(text, _):
  global out
  if(isinstance(out, outputter.quarto_outputter)):
    
    out.put_para("::: {.notes}")
    out.put_para(text)
    out.put_para(":::")
  else:
    out.put_para('---')
    out.put_para(text)

  return []


def process_picture(shape, slide_idx):
  notes = []
  if g.disable_image:
    return notes
  global picture_count
  global out

  pic_name = g.file_prefix + str(picture_count)
  pic_ext = shape.image.ext
  if not os.path.exists(g.img_path):
    os.makedirs(g.img_path)

  output_path = g.path_name_ext(g.img_path, pic_name, pic_ext)
  common_path = os.path.commonpath([g.out_path, g.img_path])
  img_outputter_path = os.path.relpath(output_path, common_path)
  with open(output_path, 'wb') as f:
    f.write(shape.image.blob)
    picture_count += 1

  # normal images
  if pic_ext != 'wmf':
    out.put_image(img_outputter_path, g.max_img_width)
    return notes

  # wmf images, try to convert, if failed, output as original
  try:
    try:
      Image.open(output_path).save(os.path.splitext(output_path)[0] + '.png')
      out.put_image(os.path.splitext(img_outputter_path)[0] + '.png', g.max_img_width)
      notes.append(f'Image {output_path} in slide {slide_idx} converted to png.')
    except Exception:  # Image failed, try another
      from wand.image import Image
      with Image(filename=output_path) as img:
        img.format = 'png'
        img.save(filename=os.path.splitext(output_path)[0] + '.png')
      out.put_image(os.path.splitext(img_outputter_path)[0] + '.png', g.max_img_width)
      notes.append(f'Image {output_path} in slide {slide_idx} converted to png111.')
  except Exception as e:
    notes.append(
        f'Cannot convert image {output_path} in slide {slide_idx} to png, this probably won\'t be displayed correctly.'
    )
    out.put_image(img_outputter_path, g.max_img_width)
  return notes

def process_table(shape, _):
  global out
  table = [[cell.text for cell in row.cells] for row in shape.table.rows]
  if len(table) > 0:
    out.put_table(table)
  return []

def ungroup_shapes(shapes):
  res = []
  for shape in shapes:
    try:
      if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        res.extend(ungroup_shapes(shape.shapes))
      else:
        res.append(shape)
    except Exception as e:
      print(f'failed to load shape {shape}, skipped. error: {e}')
  return res

def process_shapes(current_shapes, slide_id):
      local_notes = []
      for shape in current_shapes:
        if is_title(shape):
          local_notes += process_title(shape, slide_id + 1)
        elif is_text_block(shape):
          local_notes += process_text_block(shape, slide_id + 1)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
          try:
            local_notes += process_picture(shape, slide_id + 1)
          except AttributeError as e:
            print(f'Failed to process picture, skipped: {e}')
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
          local_notes += process_table(shape, slide_id + 1)
        else:
          try:
            ph = shape.placeholder_format
            if ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, "image") and getattr(shape, "image"):
              local_notes += process_picture(shape, slide_id + 1)
          except:
            pass

      return(local_notes)

# main
def parse(prs, outputer):
  global out
  out = outputer
  notes = []

  print("Starting conversion")

  # Adding inclusion of header for the first slide
  out.put_header()

  for idx, slide in enumerate(tqdm(prs.slides, desc='Converting slides')):
  # for idx, slide in enumerate(prs.slides):
    if g.page is not None and idx + 1 != g.page:
        continue
    shapes = []
    try:
      shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
    except:
      print('Bad shapes encountered in this slide. Please check or move them and try again.')
      print('shapes:')
      try:
        for sp in slide.shapes:
          print(sp.shape_type)
          print(sp.top, sp.left, sp.width, sp.height)
      except:
        print('failed to print all bad shapes.')

    process_shapes(shapes, idx + 1)

    if not g.disable_notes and slide.has_notes_slide:
      text = slide.notes_slide.notes_text_frame.text
      if text:
        notes += process_notes(text, idx + 1)
    if idx < len(prs.slides)-1 and g.enable_slides:
        out.put_para("\n---\n")
  out.close()

  if len(notes) > 0:
    print('Process finished with notice:')
    for note in notes:
      print(note)


# Alternative parse function for qmd files
def parse_alt(prs, outputer):
  global out
  out = outputer
  notes = []

  slide_width = pptx.util.Length(prs.slide_width)
  slide_width_mm = slide_width.mm

  t_vector = np.arange(1, slide_width_mm)

  print("Starting convertion to qmd file")

  # Adding inclusion of header for the first slide
  out.put_header()

  for idx, slide in enumerate(tqdm(prs.slides, desc='Converting slides')):
  # for idx, slide in enumerate(prs.slides):
    if g.page is not None and idx + 1 != g.page:
        continue
    shapes = []
    try:
      shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
    except:
      print('Bad shapes encountered in this slide. Please check or move them and try again.')
      print('shapes:')
      try:
        for sp in slide.shapes:
          print(sp.shape_type)
          print(sp.top, sp.left, sp.width, sp.height)
      except:
        print('failed to print all bad shapes.')

    pdf_modelo = is_two_column_text(slide)
    
    if(pdf_modelo):
      # Model to infer number of columns

      salida = map(lambda mu, sigma: normal_pdf(t_vector, mu, sigma), pdf_modelo[0], pdf_modelo[1])
      sum_of_gaussian = np.mean(list(salida), axis=0)
      parameters = fit_column_model(t_vector, sum_of_gaussian)

      num_cols = int(len(parameters)/2)

      dict_shapes = assign_shapes(slide, parameters, num_cols, slide_width_mm=slide_width_mm)
      print("Cantidad de elementos")
      print("pre: %d"%len(dict_shapes["shapes_pre"]))
      print("l: %d"%len(dict_shapes["shapes_l"]))
      print("c: %d"%len(dict_shapes["shapes_c"]))
      print("r: %d"%len(dict_shapes["shapes_r"]))
            
      notes.extend(process_shapes(dict_shapes["shapes_pre"], idx))
      if num_cols == 1:
        pass
      
      elif num_cols == 2:
        out.put_para(':::: {.columns}')

        out.put_para('::: {.column width="50%"}')
        notes.extend(process_shapes(dict_shapes["shapes_l"], idx))
        out.put_para(':::')
        out.put_para('::: {.column width="50%"}')
        notes.extend(process_shapes(dict_shapes["shapes_r"], idx))
        out.put_para(':::')
        out.put_para('::::')


      elif num_cols == 3:
        out.put_para(':::: {.columns}')
        
        if(dict_shapes["shapes_l"]):
          out.put_para('::: {.column width="33%"}')
          notes.extend(process_shapes(dict_shapes["shapes_l"], idx))
          out.put_para(':::')

        if(dict_shapes["shapes_c"]):
          out.put_para('::: {.column width="33%"}')
          notes.extend(process_shapes(dict_shapes["shapes_c"], idx))
          out.put_para(':::')
        
        if(dict_shapes["shapes_r"]):
          out.put_para('::: {.column width="33%"}')
          notes.extend(process_shapes(dict_shapes["shapes_r"], idx))
          out.put_para(':::')
        
        out.put_para('::::')

      else:
        raise(ValueError, "Number of columns not allowed")

    else:
      notes.extend(process_shapes(shapes, idx))


    if not g.disable_notes and slide.has_notes_slide:
      text = slide.notes_slide.notes_text_frame.text
      if text:
        notes += process_notes(text, idx + 1)
    if idx < len(prs.slides)-1 and g.enable_slides:
        out.put_para("\n---\n")
  out.close()

  if len(notes) > 0:
    print('Process finished with notice:')
    for note in notes:
      print(note)