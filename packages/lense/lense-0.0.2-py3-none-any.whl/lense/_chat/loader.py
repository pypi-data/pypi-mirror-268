
# pip install PyPDF2
# pip install python-docx
# pip install python-pptx
# pip install llama_index
import warnings
warnings.filterwarnings('ignore')
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pathlib
import warnings
from llama_index.core import SimpleDirectoryReader, Document
from llama_index.readers.file import PyMuPDFReader
from lense._chat.error_handling import *
from lense._chat.replace_file import *

def load_csv(file_path):

    try:
        data = pd.read_csv(file_path)
        error_score = error_handling(data)
        print("error_score:",error_score)
        if error_score>0:
            choice = input("Would you rather want to fix the errors if any (Yes/Y or No/N): ")
            if 'y' in choice.lower() :
                return fix_errors(data)
            else:
                return data
        else:
            print("the provided data is well structured and has no missing values in it")
            return data
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def load_excel(file_path):
    try:
        return pd.read_excel(file_path)
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def load_pdf(file_path):
    try:
        loader = PyMuPDFReader()
        documents = loader.load(file_path=file_path)
        return documents
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def load_docx(file_path):
    try:
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        if text:
            document = Document(text=text)
            print("Document loaded successfully.")
            return document
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def load_ppt(file_path):
    try:
        presentation = Presentation(file_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        if text:
            document = Document(text=text)
            print("Document loaded successfully.")
            return document
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def load_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = ''
            text += file.read()
        if text:
            document = Document(text=text)
            print("Document loaded successfully.")
            return document
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None   

def load_document(filepath):
  """
  Loads a document based on its file extension.

  Args:
      filepath (str): Path to the document file.

  Returns:
      object: The loaded document content or None if format unsupported.
  """
  filepath = pathlib.Path(filepath)
  if not filepath.is_file():
    raise FileNotFoundError(f"File not found: {filepath}")

  # Try specific libraries based on extension (update as needed)
  extension = filepath.suffix.lower()
  if extension == ".txt":
    return load_txt(filepath)
  elif extension in (".docx", ".docm"):
    return load_docx(filepath)
  elif extension == ".pdf":
    return load_pdf(filepath)
  elif extension == ".csv":
    return load_csv(filepath)
  elif extension == ".xlsx":
    return load_excel(filepath)
  elif extension == ".ppt":
    return load_ppt(filepath)

  # Unknown format
  warnings.warn(f"Unsupported file format: {extension}")
  return None

def load_folder(folder_path):
    try:
        reader = SimpleDirectoryReader(folder_path)
        documents = reader.load_data()
        return documents
    except Exception as e:
        print(f"Error reading folder '{folder_path}': {e}")
        return None

