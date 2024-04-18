import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        pdf = PdfReader(f)
        for page in pdf.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    ppt = Presentation(file_path)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")
